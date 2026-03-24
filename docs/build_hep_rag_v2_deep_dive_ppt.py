from __future__ import annotations

import sqlite3
import sys
import textwrap
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from hep_rag_v2.cli import _audit_document_payload  # noqa: E402
from hep_rag_v2.db import connect  # noqa: E402
from hep_rag_v2.graph import graph_neighbors  # noqa: E402
from hep_rag_v2.search import search_assets_bm25, search_chunks_bm25  # noqa: E402
from hep_rag_v2.textnorm import normalize_display_text  # noqa: E402


OUT_PATH = ROOT / "docs" / "hep-rag-v2-deep-dive-cn.pptx"
EXAMPLE_WORK_ID = 3
EXAMPLE_DOCUMENT_ID = 303
READY_DOCS = 87
TOTAL_FULLTEXT_DOCS = 100

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"

NAVY = RGBColor(24, 54, 93)
BLUE = RGBColor(42, 110, 187)
TEAL = RGBColor(21, 145, 123)
GREEN = RGBColor(82, 163, 94)
AMBER = RGBColor(214, 145, 29)
RED = RGBColor(185, 64, 64)
INK = RGBColor(35, 43, 60)
MUTED = RGBColor(96, 106, 122)
BG = RGBColor(246, 248, 251)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(210, 216, 226)
SOFT_BLUE = RGBColor(233, 241, 251)
SOFT_GREEN = RGBColor(234, 246, 241)
SOFT_AMBER = RGBColor(252, 245, 229)
SOFT_RED = RGBColor(252, 238, 238)
SOFT_GRAY = RGBColor(238, 242, 247)

BLOCK_TYPE_LABELS = {
    "text": "文字",
    "image": "图片",
    "equation": "公式",
    "table": "表格",
}
BLOCK_ROLE_LABELS = {
    "title": "标题",
    "front_matter": "封面信息",
    "abstract": "摘要",
    "body": "正文",
    "bibliography": "参考文献",
    "back_matter": "附录后记",
}
CHUNK_ROLE_LABELS = {
    "abstract_chunk": "摘要片段",
    "section_parent": "整节概览",
    "section_child": "局部正文片段",
    "asset_window": "图片表格上下文",
    "formula_window": "公式上下文",
}


def fmt_num(value: int) -> str:
    return f"{int(value):,}"


def shorten(text: str, limit: int) -> str:
    cleaned = " ".join((text or "").strip().split())
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1].rstrip() + "…"


def wrap_lines(text: str, width: int) -> str:
    return "\n".join(textwrap.wrap(text, width=width, break_long_words=False, break_on_hyphens=False))


def make_excerpt(text: str, limit: int = 220) -> str:
    return wrap_lines(shorten(text, limit), 44)


def label(mapping: dict[str, str], key: object) -> str:
    return mapping.get(str(key), str(key))


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_cell_fill(cell, color: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def style_text_frame(text_frame, margin: float = 0.08) -> None:
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    font_name: str = FONT_CN,
    align=PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
    line_color: RGBColor | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    tf = box.text_frame
    style_text_frame(tf)
    tf.clear()
    for idx, paragraph_text in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = paragraph_text
        p.alignment = align
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = paragraph_text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_reference_footer(slide, text: str) -> None:
    add_textbox(
        slide,
        0.6,
        7.12,
        12.0,
        0.22,
        f"参考：{text}",
        font_size=8,
        color=MUTED,
    )


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.55, 0.28, 10.6, 0.6, title, font_size=28, bold=True, color=NAVY)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.92), Inches(12.1), Inches(0.04))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.56, 1.0, 11.6, 0.35, subtitle, font_size=11, color=MUTED)


def add_panel(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    title: str | None = None,
    body: str | None = None,
    fill: RGBColor = PANEL,
    line_color: RGBColor = LINE,
    title_color: RGBColor = NAVY,
    body_color: RGBColor = INK,
    body_size: int = 13,
    title_size: int = 16,
    font_name: str = FONT_CN,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    tf = shape.text_frame
    style_text_frame(tf, margin=0.12)
    tf.clear()
    if title:
        p = tf.paragraphs[0]
        p.text = title
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = title
        run.font.name = font_name
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    if body:
        for line in str(body).split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            if p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = line
            run.font.name = font_name
            run.font.size = Pt(body_size)
            run.font.color.rgb = body_color


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: Iterable[str],
    *,
    fill: RGBColor = PANEL,
    line_color: RGBColor = LINE,
    title_color: RGBColor = NAVY,
    bullet_color: RGBColor = INK,
    bullet_size: int = 14,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    tf = shape.text_frame
    style_text_frame(tf, margin=0.12)
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = title
    run.font.name = FONT_CN
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = title_color

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = bullet
        run.font.name = FONT_CN
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color


def add_metric_card(slide, left: float, top: float, width: float, height: float, label: str, value: str, *, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    style_text_frame(tf, margin=0.08)
    tf.clear()

    p = tf.paragraphs[0]
    p.text = value
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = value
    run.font.name = FONT_CN
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    p = tf.add_paragraph()
    p.text = label
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = label
    run.font.name = FONT_CN
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED


def add_step_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: str,
    *,
    fill: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    style_text_frame(tf, margin=0.1)
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = title
    run.font.name = FONT_CN
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = NAVY

    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        run.font.name = FONT_CN
        run.font.size = Pt(11)
        run.font.color.rgb = INK


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = BLUE) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(2.2)


def add_code_card(slide, left: float, top: float, width: float, height: float, title: str, text: str, *, fill: RGBColor) -> None:
    add_panel(
        slide,
        left,
        top,
        width,
        height,
        title=title,
        body=wrap_lines(text, 36),
        fill=fill,
        body_size=11,
        font_name=FONT_MONO,
    )


def style_table(table, *, header_fill: RGBColor = NAVY, header_font_size: int = 12, body_font_size: int = 11) -> None:
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            tf = cell.text_frame
            style_text_frame(tf, margin=0.05)
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(header_font_size if row_idx == 0 else body_font_size)
                    run.font.color.rgb = PANEL if row_idx == 0 else INK
                    run.font.bold = row_idx == 0
            if row_idx == 0:
                set_cell_fill(cell, header_fill)
            else:
                set_cell_fill(cell, PANEL)


def load_data() -> dict[str, object]:
    with connect() as conn:
        snapshot = dict(
            conn.execute(
                """
                SELECT
                  (SELECT COUNT(*) FROM works) AS works,
                  (SELECT COUNT(*) FROM documents) AS documents,
                  (SELECT COUNT(*) FROM chunks) AS chunks,
                  (SELECT COUNT(*) FROM formulas) AS formulas,
                  (SELECT COUNT(*) FROM assets) AS assets,
                  (SELECT COUNT(*) FROM citations) AS citations,
                  (SELECT COUNT(*) FROM citations WHERE dst_work_id IS NOT NULL) AS resolved_citations,
                  (SELECT COUNT(*) FROM bibliographic_coupling_edges) AS bibliographic_coupling_edges,
                  (SELECT COUNT(*) FROM co_citation_edges) AS co_citation_edges,
                  (SELECT COUNT(*) FROM work_search) AS work_search,
                  (SELECT COUNT(*) FROM chunk_search) AS chunk_search,
                  (SELECT COUNT(*) FROM formula_search) AS formula_search,
                  (SELECT COUNT(*) FROM asset_search) AS asset_search
                """
            ).fetchone()
        )

        work = conn.execute(
            """
            SELECT work_id, canonical_source, canonical_id, title, abstract, year
            FROM works
            WHERE work_id = ?
            """,
            (EXAMPLE_WORK_ID,),
        ).fetchone()
        document = conn.execute(
            """
            SELECT document_id, work_id, parser_name, parser_version, parse_status, parsed_dir, manifest_path
            FROM documents
            WHERE document_id = ?
            """,
            (EXAMPLE_DOCUMENT_ID,),
        ).fetchone()

        section_count = int(
            conn.execute("SELECT COUNT(*) FROM document_sections WHERE document_id = ?", (EXAMPLE_DOCUMENT_ID,)).fetchone()[0]
        )
        block_count = int(conn.execute("SELECT COUNT(*) FROM blocks WHERE document_id = ?", (EXAMPLE_DOCUMENT_ID,)).fetchone()[0])
        chunk_count = int(conn.execute("SELECT COUNT(*) FROM chunks WHERE document_id = ?", (EXAMPLE_DOCUMENT_ID,)).fetchone()[0])
        formula_count = int(conn.execute("SELECT COUNT(*) FROM formulas WHERE document_id = ?", (EXAMPLE_DOCUMENT_ID,)).fetchone()[0])
        asset_count = int(conn.execute("SELECT COUNT(*) FROM assets WHERE document_id = ?", (EXAMPLE_DOCUMENT_ID,)).fetchone()[0])
        outgoing_citations = int(conn.execute("SELECT COUNT(*) FROM citations WHERE src_work_id = ?", (EXAMPLE_WORK_ID,)).fetchone()[0])
        incoming_citations = int(conn.execute("SELECT COUNT(*) FROM citations WHERE dst_work_id = ?", (EXAMPLE_WORK_ID,)).fetchone()[0])

        block_roles = {
            str(row["block_role"]): int(row["n"])
            for row in conn.execute(
                """
                SELECT block_role, COUNT(*) AS n
                FROM blocks
                WHERE document_id = ?
                GROUP BY block_role
                ORDER BY block_role
                """,
                (EXAMPLE_DOCUMENT_ID,),
            ).fetchall()
        }
        chunk_roles = {
            str(row["chunk_role"]): int(row["n"])
            for row in conn.execute(
                """
                SELECT chunk_role, COUNT(*) AS n
                FROM chunks
                WHERE document_id = ?
                GROUP BY chunk_role
                ORDER BY chunk_role
                """,
                (EXAMPLE_DOCUMENT_ID,),
            ).fetchall()
        }

        sample_blocks = [
            dict(row)
            for row in conn.execute(
                """
                SELECT block_id, block_type, page, block_role, is_heading, clean_text
                FROM blocks
                WHERE document_id = ?
                  AND block_id IN (57969, 57970, 57972, 58013)
                ORDER BY block_id
                """,
                (EXAMPLE_DOCUMENT_ID,),
            ).fetchall()
        ]
        raw_clean = dict(
            conn.execute(
                """
                SELECT block_id, raw_text, clean_text
                FROM blocks
                WHERE block_id = 57970
                """
            ).fetchone()
        )
        raw_orphan = dict(
            conn.execute(
                """
                SELECT block_id, raw_text, clean_text
                FROM blocks
                WHERE block_id = 57980
                """
            ).fetchone()
        )
        equation_block = dict(
            conn.execute(
                """
                SELECT block_id, raw_text, clean_text
                FROM blocks
                WHERE block_id = 58013
                """
            ).fetchone()
        )
        equation_chunk = dict(
            conn.execute(
                """
                SELECT chunk_id, chunk_role, clean_text
                FROM chunks
                WHERE chunk_id = 11261
                """
            ).fetchone()
        )
        chunk_examples = [
            dict(row)
            for row in conn.execute(
                """
                SELECT chunk_id, chunk_role, page_hint, clean_text
                FROM chunks
                WHERE chunk_id IN (11239, 11240, 11243, 11261)
                ORDER BY chunk_id
                """
            ).fetchall()
        ]
        figure_row = dict(
            conn.execute(
                """
                SELECT a.asset_id, a.asset_path, a.caption, b.clean_text
                FROM assets a
                LEFT JOIN blocks b ON b.block_id = a.block_id
                WHERE a.document_id = ?
                ORDER BY a.asset_id
                LIMIT 1
                """,
                (EXAMPLE_DOCUMENT_ID,),
            ).fetchone()
        )
        image_path = Path(str(document["parsed_dir"])) / "raw" / str(figure_row["asset_path"])

        audit = _audit_document_payload(conn, document_id=EXAMPLE_DOCUMENT_ID, limit=6)
        query = "four photons higgs pseudoscalar"
        display_query = "四光子、Higgs、轻赝标量"
        chunk_hits = search_chunks_bm25(conn, query=query, collection="cms_rare_decay", limit=3)
        asset_hits = search_assets_bm25(conn, query="feynman diagram higgs pseudoscalar photons", collection="cms_rare_decay", limit=2)
        neighbors = graph_neighbors(conn, work_id=EXAMPLE_WORK_ID, edge_kind="all", collection="cms_rare_decay", limit=5)

    return {
        "snapshot": snapshot,
        "work": dict(work),
        "document": dict(document),
        "section_count": section_count,
        "block_count": block_count,
        "chunk_count": chunk_count,
        "formula_count": formula_count,
        "asset_count": asset_count,
        "outgoing_citations": outgoing_citations,
        "incoming_citations": incoming_citations,
        "block_roles": block_roles,
        "chunk_roles": chunk_roles,
        "sample_blocks": sample_blocks,
        "raw_clean": raw_clean,
        "raw_orphan": raw_orphan,
        "equation_block": equation_block,
        "equation_chunk": equation_chunk,
        "chunk_examples": chunk_examples,
        "figure": figure_row,
        "image_path": image_path,
        "audit": audit,
        "query": query,
        "display_query": display_query,
        "chunk_hits": chunk_hits,
        "asset_hits": asset_hits,
        "neighbors": neighbors,
    }


def build_slide_1(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.4), Inches(12.35), Inches(1.05))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()
    add_textbox(slide, 0.75, 0.6, 11.2, 0.45, "hep-rag-v2：把论文资料整理成可查、可串联、可写作的基础设施", font_size=28, bold=True, color=PANEL)
    add_textbox(slide, 0.78, 1.03, 10.2, 0.28, "面向高能物理论文整理、线索追踪和写作准备", font_size=12, color=RGBColor(225, 232, 243))

    add_panel(
        slide,
        0.6,
        1.75,
        7.2,
        2.1,
        title="一句话结论",
        body="现在已经能把论文的题录信息、正文片段、图片、公式和引用关系整理进同一套库里。\n底层整理过程主要靠规则和算法完成；大模型放在后面的写作和整理环节，而不是放在入库阶段。",
        fill=PANEL,
        body_size=15,
    )

    metrics = [
        ("论文总数", fmt_num(data["snapshot"]["works"]), SOFT_BLUE),
        ("已整理全文", fmt_num(data["snapshot"]["documents"]), SOFT_GREEN),
        ("可检索片段", fmt_num(data["snapshot"]["chunks"]), SOFT_AMBER),
        ("引用条目", fmt_num(data["snapshot"]["citations"]), SOFT_RED),
    ]
    x = 8.15
    for idx, (label, value, fill) in enumerate(metrics):
        add_metric_card(slide, x + idx * 1.08, 1.88, 1.0, 1.25, label, value, fill=fill)

    add_panel(
        slide,
        0.6,
        4.15,
        12.1,
        1.0,
        title="当前边界",
        body="默认不让大模型参与入库、清洗、分段和建立论文关系。大模型主要留给后面的检索辅助、证据整理和写作辅助。",
        fill=SOFT_BLUE,
        body_size=14,
    )

    add_metric_card(slide, 0.8, 5.45, 2.4, 1.2, "可进入下一步", f"{READY_DOCS}/{TOTAL_FULLTEXT_DOCS}", fill=SOFT_GREEN)
    add_metric_card(slide, 3.4, 5.45, 2.5, 1.2, "已连上的引用", fmt_num(data["snapshot"]["resolved_citations"]), fill=SOFT_BLUE)
    add_metric_card(slide, 6.15, 5.45, 2.8, 1.2, "共享参考文献关系", fmt_num(data["snapshot"]["bibliographic_coupling_edges"]), fill=SOFT_AMBER)
    add_metric_card(slide, 9.25, 5.45, 2.2, 1.2, "共同被引关系", fmt_num(data["snapshot"]["co_citation_edges"]), fill=SOFT_RED)


def build_slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "1. 先讲问题：这套东西是拿来做什么的", "先讲用途，再讲技术。")

    add_bullets(
        slide,
        0.6,
        1.45,
        5.75,
        4.65,
        "它要解决的核心问题",
        [
            "给研究者一套能持续扩量的论文资料底座",
            "支持找同主题论文、追踪引用线索、做专题综述、整理方法路线",
            "未来目标规模可以到大约 100k 篇，所以不能一开始就要求每篇都做昂贵全文处理",
            "要把“大规模收集文献信息”和“精细整理全文证据”分开处理",
        ],
        fill=PANEL,
        bullet_size=16,
    )

    add_bullets(
        slide,
        6.6,
        1.45,
        6.15,
        4.65,
        "为什么不能只堆 PDF 然后搜关键词",
        [
            "标题、作者、正文、图注、公式、参考文献会混在一起",
            "PDF/LaTeX 解析有大量噪声：空格断裂、符号拆散、公式命令残留",
            "检索命中后，不知道这一段是不是可作为证据引用的正文",
            "单纯关键词搜索缺少论文之间的结构关系，无法自然扩展到近邻工作",
        ],
        fill=PANEL,
        bullet_size=16,
    )

    add_step_box(slide, 0.9, 6.25, 3.25, 0.75, "路线 A：先铺开文献信息", "先整理标题、摘要、作者和参考文献", fill=SOFT_BLUE)
    add_step_box(slide, 4.55, 6.25, 3.25, 0.75, "路线 B：全文只做重点子集", "把正文、图片、公式整理成可检索证据", fill=SOFT_GREEN)
    add_step_box(slide, 8.2, 6.25, 3.55, 0.75, "路线 C：写作辅助以后再接", "放在应用层，不直接改底层事实库", fill=SOFT_AMBER)


def build_slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "2. 一页看完整件事", "前 5 步不依赖大模型。")

    steps = [
        ("收集输入", "文献信息\n全文解析结果", SOFT_BLUE),
        ("整理论文卡片", "编号、题目、作者\n摘要、参考文献", SOFT_GREEN),
        ("整理全文结构", "章节、页面片段\n图片、公式", SOFT_AMBER),
        ("清洗文字", "去引用残留\n修解析噪声", SOFT_RED),
        ("生成检索片段", "整节概览\n局部证据片段", SOFT_BLUE),
        ("补上论文关系", "相近论文\n引用网络", SOFT_GREEN),
    ]
    x = 0.45
    y = 2.3
    w = 1.9
    for idx, (title, body, fill) in enumerate(steps):
        add_step_box(slide, x + idx * 2.1, y, w, 1.45, title, body, fill=fill)
        if idx < len(steps) - 1:
            add_arrow(slide, x + idx * 2.1 + w, y + 0.73, x + (idx + 1) * 2.1, y + 0.73)

    add_bullets(
        slide,
        0.7,
        4.55,
        4.1,
        2.0,
        "当前已经落地的部分",
        [
            "文献信息入库",
            "全文结构整理",
            "规则化文字清洗",
            "关键词检索加论文关系扩展",
        ],
        fill=PANEL,
        bullet_size=14,
    )
    add_bullets(
        slide,
        4.95,
        4.55,
        3.9,
        2.0,
        "当前明确不做的事",
        [
            "不用大模型自动切块",
            "不用大模型自动建图",
            "不让模型直接改原始事实库",
            "不把参考文献正文混进主检索流",
        ],
        fill=PANEL,
        bullet_size=14,
    )
    add_bullets(
        slide,
        9.0,
        4.55,
        3.6,
        2.0,
        "下一阶段再加的能力",
        [
            "向量表示",
            "检索意图整理",
            "证据压缩",
            "综述和写作辅助",
        ],
        fill=PANEL,
        bullet_size=14,
    )


def build_slide_4(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "3. 用一篇真实论文把故事讲清楚", "所有层级和数字都来自当前数据库。")

    work_title = normalize_display_text(str(data["work"]["title"]))
    short_title = "例子：一篇研究 H -> aa -> 4γ 的 CMS 论文"
    add_panel(
        slide,
        0.55,
        1.35,
        7.0,
        1.15,
        title=short_title,
        body=wrap_lines(work_title, 70),
        fill=PANEL,
        body_size=12,
    )

    stats = [
        ("年份", str(data["work"]["year"]), SOFT_BLUE),
        ("章节数", str(data["section_count"]), SOFT_GREEN),
        ("页面片段", str(data["block_count"]), SOFT_AMBER),
        ("检索片段", str(data["chunk_count"]), SOFT_RED),
        ("图片表格", str(data["asset_count"]), SOFT_BLUE),
        ("参考文献", str(data["outgoing_citations"]), SOFT_GREEN),
    ]
    for idx, (label, value, fill) in enumerate(stats):
        row = idx // 3
        col = idx % 3
        add_metric_card(slide, 0.7 + col * 1.55, 2.8 + row * 1.25, 1.35, 1.05, label, value, fill=fill)

    add_panel(
        slide,
        0.58,
        5.48,
        6.95,
        1.35,
        title="摘要里已经能看到未来检索的“证据句”",
        body=make_excerpt(str(data["work"]["abstract"]), 330),
        fill=PANEL,
        body_size=11,
    )

    if Path(data["image_path"]).exists():
        slide.shapes.add_picture(str(data["image_path"]), Inches(7.9), Inches(1.55), width=Inches(4.3))
    add_panel(
        slide,
        7.65,
        5.48,
        4.7,
        1.35,
        title="这张图为什么重要",
        body="它不会孤零零留在图片文件夹里，而是和前后正文一起进入库中。这样后面既能按图片找，也能把相关说明一起带出来。",
        fill=PANEL,
        body_size=11,
    )


def build_slide_5(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "4. 这篇论文在库里会经历哪些层", "从文献卡片，到可检索片段，再到论文关系。")

    top_y = 1.65
    add_step_box(slide, 0.45, top_y, 1.75, 1.15, "文献卡片", "题目、摘要、年份\n参考文献", fill=SOFT_BLUE)
    add_step_box(slide, 2.45, top_y, 1.75, 1.15, "全文记录", "一份可处理的\n论文全文", fill=SOFT_GREEN)
    add_step_box(slide, 4.45, top_y, 1.75, 1.15, "章节", f"{data['section_count']} 个章节\n摘要、引言等", fill=SOFT_AMBER)
    add_step_box(slide, 6.45, top_y, 1.75, 1.15, "页面片段", f"{data['block_count']} 个小片段\n文字、图片、公式", fill=SOFT_RED)
    add_step_box(slide, 8.45, top_y, 1.75, 1.15, "检索片段", f"{data['chunk_count']} 个可用片段\n按用途重组", fill=SOFT_BLUE)
    add_step_box(slide, 10.45, top_y, 2.25, 1.15, "论文关系", f"{data['outgoing_citations']} 条引用线索\n连接相关论文", fill=SOFT_GREEN)

    for x in (2.2, 4.2, 6.2, 8.2, 10.2):
        add_arrow(slide, x, top_y + 0.57, x + 0.25, top_y + 0.57)

    add_panel(
        slide,
        0.55,
        3.25,
        3.95,
        2.6,
        title="页面片段层的真实内容",
        body="\n".join(
            [
                f"{row['block_id']} | {label(BLOCK_TYPE_LABELS, row['block_type'])} | 第 {row['page']} 页 | {shorten(str(row['clean_text']), 52)}"
                for row in data["sample_blocks"]
            ]
        ),
        fill=PANEL,
        body_size=11,
    )

    chunk_role_text = "\n".join(
        [
            f"{label(CHUNK_ROLE_LABELS, role)}：{count}"
            for role, count in data["chunk_roles"].items()
        ]
    )
    add_panel(
        slide,
        4.7,
        3.25,
        2.55,
        2.6,
        title="可检索片段的几种类型",
        body=chunk_role_text,
        fill=PANEL,
        body_size=12,
    )

    add_panel(
        slide,
        7.45,
        3.25,
        5.2,
        2.6,
        title="最简单的理解",
        body="同一篇论文不会直接变成“一整条长文本”。\n它会先被拆成章节和页面片段，再按照用途重新整理成几种可检索片段：\n有的负责整节概览，有的负责局部证据，有的专门把图片或公式前后的说明保留下来。",
        fill=PANEL,
        body_size=12,
    )

    add_panel(
        slide,
        0.55,
        6.05,
        12.1,
        0.72,
        title="一句话",
        body="你可以把它理解为：一篇论文先被拆成很多有明确用途的小片段，再和其他论文通过引用关系连接起来。",
        fill=SOFT_BLUE,
        body_size=13,
    )


def build_slide_6(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "5. 第一步：先把文献信息和引用关系整理好", "这一步还没有处理全文正文。")

    add_bullets(
        slide,
        0.55,
        1.45,
        5.1,
        3.55,
        "文献信息整理的核心规则",
        [
            "先给每篇论文确定一个稳定身份，优先用 INSPIRE 编号，其次用 arXiv 和 DOI",
            "把题目、摘要、作者、合作组、期刊和主题词整理进对应表里",
            "把论文自带的参考文献一条条写进引用关系表",
            "按目标编号去重，避免同一条参考文献重复写入",
            "如果被引用的论文已经在库里，就直接把两篇论文连起来；暂时找不到的先记下来",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    add_bullets(
        slide,
        5.85,
        1.45,
        3.15,
        3.55,
        "论文关系怎么建",
        [
            "第一种：两篇论文引用了很多相同的前人工作",
            "第二种：很多后来的论文总是同时提到它们",
            "共同条目越多，两篇论文越可能属于同一条研究线",
            "小型试验库阈值可以放松；规模变大后再收紧",
        ],
        fill=PANEL,
        bullet_size=13,
    )

    table_shape = slide.shapes.add_table(6, 4, Inches(9.15), Inches(1.48), Inches(3.2), Inches(3.45))
    table = table_shape.table
    headers = ["相关论文", "关系", "共同条目", "紧密度"]
    for idx, value in enumerate(headers):
        table.cell(0, idx).text = value
    for row_idx, row in enumerate(data["neighbors"][:5], start=1):
        table.cell(row_idx, 0).text = str(row["canonical_id"])
        table.cell(row_idx, 1).text = "共参考文献" if row["edge_kind"] == "bibliographic_coupling" else "共被引"
        table.cell(row_idx, 2).text = str(row["shared_count"])
        table.cell(row_idx, 3).text = f"{float(row['score']):.2f}"
    style_table(table)

    add_panel(
        slide,
        0.55,
        5.2,
        12.0,
        1.4,
        title="这个例子里发生了什么",
        body=f"这篇论文本身带有 {data['outgoing_citations']} 条引用线索。把这些线索补齐以后，系统就能自然连到同主题的 H->aa 系列工作。对现在这 101 篇试验语料来说，这比只搜标题和摘要更适合做专题梳理。",
        fill=SOFT_GREEN,
        body_size=13,
    )


def build_slide_7(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "6. 第二步：把全文整理成章节、页面片段、图片和公式", "这里使用 MinerU 的解析结果，但后面还会再做规则判断。")

    add_bullets(
        slide,
        0.55,
        1.45,
        4.1,
        4.35,
        "全文是怎样被整理出结构的",
        [
            "逐条读取 MinerU 解析出来的内容列表",
            "先判断每一小块是文字、图片、表格还是公式",
            "再判断这块像不像标题，能不能作为新章节的开头",
            "系统会根据当前位置判断自己已经进入摘要、正文还是参考文献",
            "每遇到新标题，就在库里新建一个章节节点",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    add_bullets(
        slide,
        4.85,
        1.45,
        3.0,
        4.35,
        "保留什么，排除什么",
        [
            "标题、封面信息、参考文献和后记都会保留在库里",
            "但这些默认不进主检索流",
            "公式会单独登记成公式记录",
            "图片和表格会单独登记成图表记录",
            "正文页面片段才是后面重组检索片段的主要输入",
        ],
        fill=PANEL,
        bullet_size=13,
    )

    table_shape = slide.shapes.add_table(5, 4, Inches(8.05), Inches(1.5), Inches(4.4), Inches(4.25))
    table = table_shape.table
    for idx, head in enumerate(["片段编号", "类型", "位置", "例子"]):
        table.cell(0, idx).text = head
    for row_idx, row in enumerate(data["sample_blocks"], start=1):
        table.cell(row_idx, 0).text = str(row["block_id"])
        table.cell(row_idx, 1).text = label(BLOCK_TYPE_LABELS, row["block_type"])
        table.cell(row_idx, 2).text = label(BLOCK_ROLE_LABELS, row["block_role"])
        table.cell(row_idx, 3).text = shorten(str(row["clean_text"]), 45)
    style_table(table)

    add_panel(
        slide,
        0.55,
        6.0,
        12.05,
        0.72,
        title="关键点",
        body="这一步不是让模型去“理解全文结构”，而是把解析结果经过一组明确规则，整理成可检查、可回溯的章节树和页面片段树。",
        fill=SOFT_BLUE,
        body_size=13,
    )
    add_reference_footer(
        slide,
        "Bin Wang et al. (2024), MinerU: An Open-Source Solution for Precise Document Content Extraction; "
        "Jiawei Wang et al. (2024), Detect-Order-Construct: A Tree Construction based Approach for Hierarchical Document Structure Analysis.",
    )


def build_slide_8(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "7. 第三步：清洗文字，并自动检查质量", "这里专门处理空格断裂、引用残留和公式残留。")

    add_bullets(
        slide,
        0.55,
        1.4,
        3.9,
        5.6,
        "主要清洗规则",
        [
            "去掉 HTML 标签、控制字符和奇怪空白",
            "数字拼回去：1 0 0 -> 100",
            "字母拼回去：G e V -> GeV",
            "去掉正文里的引用残留： [1–3] / Ref. [37]",
            "把常见 LaTeX 命令改写成可读文字",
            "修正常见解析残留：centerofmass -> center-of-mass",
            "把 frac S B 这类残留整理成 (S) / (B)",
            "像 1 0 ^ { - 8 } 这样的残留，本质上是数学符号被拆散了",
        ],
        fill=PANEL,
        bullet_size=13,
    )

    raw_text = make_excerpt(str(data["raw_clean"]["raw_text"]), 255)
    clean_text = make_excerpt(str(data["raw_clean"]["clean_text"]), 255)
    add_code_card(slide, 4.7, 1.48, 3.65, 2.35, "例 1：清洗前", raw_text, fill=SOFT_RED)
    add_code_card(slide, 8.55, 1.48, 3.65, 2.35, "例 1：清洗后", clean_text, fill=SOFT_GREEN)

    orphan_text = make_excerpt(str(data["raw_orphan"]["clean_text"]), 220)
    add_code_card(slide, 4.7, 4.0, 3.65, 1.95, "例 2：当前残余问题", orphan_text, fill=SOFT_AMBER)

    eq_raw = make_excerpt(str(data["equation_block"]["raw_text"]).replace("\n", " "), 180)
    eq_chunk = make_excerpt(str(data["equation_chunk"]["clean_text"]), 230)
    add_code_card(slide, 8.55, 4.0, 3.65, 1.95, "例 3：公式整理后的样子", f"原始公式: {eq_raw}\n----\n进入正文后的样子: {eq_chunk}", fill=SOFT_BLUE)

    audit = data["audit"]
    noise_chunk_hits = audit["noise"]["retrievable_chunks"]["total_hits"]
    noise_block_hits = audit["noise"]["retrievable_blocks"]["total_hits"]
    add_panel(
        slide,
        4.68,
        6.15,
        7.55,
        0.78,
        title="自动检查是怎么做的",
        body=f"系统会自动扫描可检索正文里的残余噪声。当前这篇论文仍有 {noise_chunk_hits} 处片段级问题、{noise_block_hits} 处页面片段级问题，所以还需要继续清理。",
        fill=SOFT_BLUE,
        body_size=12,
    )


def build_slide_9(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "8. 第四步：把页面片段重组成可检索片段", "这里也回答：为什么不会被切成半句话。")

    add_step_box(slide, 0.55, 1.65, 2.0, 1.05, "输入", "可用的正文片段\n排除参考文献", fill=SOFT_BLUE)
    add_step_box(slide, 2.8, 1.65, 2.0, 1.05, "线性整理", "正文直接进入\n公式先变成简短说明", fill=SOFT_GREEN)
    add_step_box(slide, 5.05, 1.65, 2.0, 1.05, "合并", "看前后是不是同一句\n必要时拼回去", fill=SOFT_AMBER)
    add_step_box(slide, 7.3, 1.65, 2.0, 1.05, "切分", "尽量按句子切开\n过长时再细分", fill=SOFT_RED)
    add_step_box(slide, 9.55, 1.65, 2.45, 1.05, "输出", "整节概览 / 局部片段\n图表上下文 / 公式上下文", fill=SOFT_BLUE)
    for x in (2.55, 4.8, 7.05, 9.3):
        add_arrow(slide, x, 2.17, x + 0.25, 2.17)

    add_bullets(
        slide,
        0.55,
        3.05,
        4.0,
        3.45,
        "重组规则",
        [
            "每个片段默认控制在大约 2400 个字符以内",
            "如果整章不长，会先保留一份整节概览",
            "局部正文片段会带一点点重叠，避免信息在边界处断掉",
            "图片和表格会连同前后说明一起单独保留",
            "公式会保留前后解释文字，而不是把整段公式硬塞进正文",
        ],
        fill=PANEL,
        bullet_size=13,
    )

    add_bullets(
        slide,
        4.8,
        3.05,
        3.95,
        3.45,
        "为什么不会切成半句话/半个词",
        [
            "先按段落和句号切开，而不是直接按字符数硬砍",
            "如果一句话还是太长，再继续往下细分",
            "切开后的词会再重新拼成正常字符串",
            "不会直接按字符数硬砍，所以不会出现半个英文单词",
            "系统还会尽量避开 -, /, = 这类容易把式子切坏的位置",
        ],
        fill=PANEL,
        bullet_size=13,
    )

    example_chunk = next(row for row in data["chunk_examples"] if row["chunk_id"] == 11261)
    add_code_card(slide, 9.0, 3.05, 3.3, 3.45, "实际片段示例", make_excerpt(str(example_chunk["clean_text"]), 420), fill=SOFT_GREEN)
def build_slide_10(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "9. 第五步：用户检索时会发生什么", "先找命中的片段，再把图和上下文带出来。")

    add_bullets(
        slide,
        0.55,
        1.42,
        4.15,
        4.9,
        "当前检索方式",
        [
            "先把检索词整理成适合搜索的形式",
            "分别在论文卡片、正文片段、公式和图片说明里查找",
            "排序时优先把正文片段放前面，图表和公式上下文放后面",
            "如果命中的是图片或公式，还会自动带上前后说明文字",
            "当前先用传统关键词检索，向量检索以后再接",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    hit = data["chunk_hits"][0]
    add_panel(
        slide,
        4.95,
        1.45,
        4.05,
        2.15,
        title=f"示例检索词：{data['display_query']}",
        body=f"命中的正文片段：第 {hit['page_hint']} 页，类型是“{label(CHUNK_ROLE_LABELS, hit['chunk_role'])}”\n{make_excerpt(str(hit['clean_text']), 250)}",
        fill=SOFT_BLUE,
        body_size=11,
    )

    asset_hit = data["asset_hits"][0]
    add_panel(
        slide,
        9.2,
        1.45,
        3.2,
        2.15,
        title="图片命中示例",
        body=f"第 {asset_hit['page']} 页的图片说明：\n{wrap_lines(shorten(str(asset_hit['caption']), 130), 28)}",
        fill=SOFT_GREEN,
        body_size=11,
    )

    add_step_box(slide, 1.1, 4.15, 2.2, 0.95, "检索词", "用户输入的问题\n或关键词", fill=SOFT_BLUE)
    add_step_box(slide, 3.65, 4.15, 2.2, 0.95, "查找", "论文卡片\n正文片段\n公式\n图片", fill=SOFT_GREEN)
    add_step_box(slide, 6.2, 4.15, 2.4, 0.95, "排序", "先给正文证据\n再给图和公式说明", fill=SOFT_AMBER)
    add_step_box(slide, 8.95, 4.15, 2.55, 0.95, "证据包", "可引用的文字证据\n加上下文", fill=SOFT_RED)
    for x in (3.3, 5.85, 8.6):
        add_arrow(slide, x, 4.62, x + 0.25, 4.62)

    add_panel(
        slide,
        0.9,
        5.45,
        10.7,
        1.0,
        title="为什么这一步已经很有价值",
        body="因为系统已经把正文、图片说明、公式说明和参考文献分开了。即便还只是关键词检索，命中的内容也比原始 PDF 搜索更接近真正能引用的证据。",
        fill=PANEL,
        body_size=12,
    )
    add_reference_footer(
        slide,
        "Jiasheng Zhang et al. (2024), LitFM: A Retrieval Augmented Structure-aware Foundation Model For Citation Graphs; "
        "Yuntong Hu et al. (2025), CG-RAG: Research Question Answering by Citation Graph Retrieval-Augmented LLMs.",
    )


def build_slide_11(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "10. 只靠关键词还不够：还要把相关论文串起来", "找到一篇之后，还要知道它周围是谁。")

    add_bullets(
        slide,
        0.55,
        1.42,
        4.35,
        4.9,
        "两种边的直觉理解",
        [
            "第一种：两篇论文引用了很多相同的前人工作",
            "第二种：很多后来的论文总是同时提到它们",
            "前者更像“参考文献口味相似”",
            "后者更像“后人把它们看成同一簇主题”",
            "做综述时，这两种边都很有价值",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    table_shape = slide.shapes.add_table(6, 4, Inches(5.15), Inches(1.45), Inches(7.1), Inches(4.95))
    table = table_shape.table
    for idx, head in enumerate(["相关论文", "关系", "共同条目", "题目"]):
        table.cell(0, idx).text = head
    for row_idx, row in enumerate(data["neighbors"][:5], start=1):
        table.cell(row_idx, 0).text = str(row["canonical_id"])
        table.cell(row_idx, 1).text = "共参考文献" if row["edge_kind"] == "bibliographic_coupling" else "共被引"
        table.cell(row_idx, 2).text = str(row["shared_count"])
        table.cell(row_idx, 3).text = shorten(str(row["title"]), 62)
    style_table(table)

    add_panel(
        slide,
        0.78,
        6.05,
        11.8,
        0.72,
        title="更直观的理解",
        body="检索先找到一篇明显相关的论文，论文关系层再把它周围主题相近的论文一起捞出来，于是专题梳理不再只靠关键词穷举。",
        fill=SOFT_BLUE,
        body_size=13,
    )
    add_reference_footer(
        slide,
        "Darren Edge et al. (2024), From Local to Global: A Graph RAG Approach to Query-Focused Summarization; "
        "Haoyu Han et al. (2025), Retrieval-Augmented Generation with Graphs (GraphRAG); "
        "Yilin Xiao et al. (2025), GraphRAG-Bench: Challenging Domain-Specific Reasoning for Evaluating Graph Retrieval-Augmented Generation.",
    )


def build_slide_12(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "11. 哪些部分完全靠规则，哪些部分以后再接大模型", "把边界讲清楚，后面系统才容易维护。")

    left = slide.shapes.add_table(7, 3, Inches(0.55), Inches(1.45), Inches(6.15), Inches(4.95)).table
    for idx, head in enumerate(["步骤", "当前是否用大模型", "当前输出"]):
        left.cell(0, idx).text = head
    algorithm_rows = [
        ("整理文献信息", "否", "论文卡片、作者、引用关系"),
        ("整理全文结构", "否", "章节、页面片段、图片、公式"),
        ("清洗文字", "否", "更干净的正文和清洗标记"),
        ("重组检索片段", "否", "整节概览、局部片段、图表上下文"),
        ("检索和论文关系", "否", "命中的证据和相近论文"),
        ("自动质量检查", "否", "是否通过当前质量门槛"),
    ]
    for row_idx, row in enumerate(algorithm_rows, start=1):
        for col_idx, value in enumerate(row):
            left.cell(row_idx, col_idx).text = value
    style_table(left)

    right = slide.shapes.add_table(5, 3, Inches(7.0), Inches(1.45), Inches(5.65), Inches(4.95)).table
    for idx, head in enumerate(["以后可接的大模型环节", "输入", "预期输出"]):
        right.cell(0, idx).text = head
    llm_rows = [
        ("检索意图整理", "用户问题和范围", "一组检索关键词和筛选范围"),
        ("证据压缩", "命中的正文片段和相关论文", "按主题整理好的证据卡片"),
        ("综述初稿", "证据卡片和论文清单", "按章节展开的初稿"),
        ("思路梳理", "相关工作矩阵", "研究空档、对比和候选方向"),
    ]
    for row_idx, row in enumerate(llm_rows, start=1):
        for col_idx, value in enumerate(row):
            right.cell(row_idx, col_idx).text = value
    style_table(right, header_fill=TEAL)

    add_panel(
        slide,
        0.7,
        6.08,
        11.9,
        0.72,
        title="原则",
        body="大模型负责读证据、写总结、产出候选草稿；但不直接决定底层入库、清洗和建立论文关系的逻辑。",
        fill=SOFT_GREEN,
        body_size=13,
    )


def build_slide_13(prs: Presentation, data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "12. 现在做到哪一步了，下一步该怎么走", "先把底座讲清楚，再谈后续扩量。")

    add_metric_card(slide, 0.7, 1.55, 1.6, 1.15, "论文数", fmt_num(data["snapshot"]["works"]), fill=SOFT_BLUE)
    add_metric_card(slide, 2.5, 1.55, 1.7, 1.15, "全文数", fmt_num(data["snapshot"]["documents"]), fill=SOFT_GREEN)
    add_metric_card(slide, 4.45, 1.55, 1.6, 1.15, "检索片段", fmt_num(data["snapshot"]["chunks"]), fill=SOFT_AMBER)
    add_metric_card(slide, 6.25, 1.55, 1.6, 1.15, "图片表格", fmt_num(data["snapshot"]["assets"]), fill=SOFT_RED)
    add_metric_card(slide, 8.05, 1.55, 1.75, 1.15, "公式", fmt_num(data["snapshot"]["formulas"]), fill=SOFT_BLUE)
    add_metric_card(slide, 10.05, 1.55, 1.85, 1.15, "可继续推进", f"{READY_DOCS}/{TOTAL_FULLTEXT_DOCS}", fill=SOFT_GREEN)

    add_bullets(
        slide,
        0.55,
        3.0,
        5.35,
        3.35,
        "现在已经可以说它是什么",
        [
            "它已经不是一个只会解析 PDF 的小工具",
            "它已经能把文献信息、全文片段、图片、公式和引用关系放进同一套库里",
            "它已经有稳定的文字清洗和片段重组流程",
            "它已经能支持检索和专题梳理",
            "它已经能顺着引用关系找到相关论文",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    add_bullets(
        slide,
        6.1,
        3.0,
        6.2,
        3.35,
        "最合理的下一步",
        [
            "先把文献信息层扩到更大规模，不要一开始就上 100k 全文",
            "全文层继续只做高价值子集，先把证据检索跑顺",
            "之后再加更细的相似度和向量层",
            "最后再把综述写作和研究辅助接到大模型应用层",
            "对外展示时，核心信息是：底层论文资料库已经成形，后面还能继续长能力",
        ],
        fill=PANEL,
        bullet_size=14,
    )

    add_panel(
        slide,
        0.75,
        6.5,
        11.65,
        0.6,
        title="建议的一句话路线",
        body="先做“可扩量的文献信息层 + 小规模高质量全文证据层”，再做向量层和大模型应用层。",
        fill=SOFT_BLUE,
        body_size=13,
    )


def build_deck(data: dict[str, object]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs, data)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs, data)
    build_slide_5(prs, data)
    build_slide_6(prs, data)
    build_slide_7(prs, data)
    build_slide_8(prs, data)
    build_slide_9(prs, data)
    build_slide_10(prs, data)
    build_slide_11(prs, data)
    build_slide_12(prs)
    build_slide_13(prs, data)
    return prs


def main() -> None:
    data = load_data()
    prs = build_deck(data)
    prs.save(str(OUT_PATH))
    print(OUT_PATH)


if __name__ == "__main__":
    main()
